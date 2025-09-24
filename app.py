# app.py: A Streamlit web application for universal file-to-text conversion.
# This app allows users to upload a file (Word, Excel, PowerPoint, HTML, ZIP),
# see a preview of the extracted text, and download the full content as a .txt file.

import streamlit as st
import os
import zipfile
import io
import docx
import openpyxl
import pptx
from markdownify import markdownify as md

# --- File Processing Functions (adapted from the Colab version) ---

def convert_docx_to_text(file_stream):
    """Extracts text from a .docx file stream."""
    doc = docx.Document(file_stream)
    return "\n".join([para.text for para in doc.paragraphs])

def convert_xlsx_to_text(file_stream):
    """Extracts text from an .xlsx file stream, sheet by sheet."""
    workbook = openpyxl.load_workbook(file_stream)
    full_text = []
    for sheet_name in workbook.sheetnames:
        full_text.append(f"--- Sheet: {sheet_name} ---\n")
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows():
            row_text = "\t".join([str(cell.value) if cell.value is not None else "" for cell in row])
            full_text.append(row_text)
    return "\n".join(full_text)

def convert_pptx_to_text(file_stream):
    """Extracts text from a .pptx file stream from all shapes on all slides."""
    presentation = pptx.Presentation(file_stream)
    full_text = []
    for i, slide in enumerate(presentation.slides):
        full_text.append(f"--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text)

def convert_html_to_markdown(file_stream):
    """Converts HTML file content into Markdown using markdownify."""
    html_content = file_stream.read().decode('utf-8')
    # Using markdownify to convert HTML to Markdown text
    return md(html_content)

def process_zip_file(file_stream):
    """Extracts files from a .zip archive in memory and converts them."""
    full_text = []
    with zipfile.ZipFile(file_stream, 'r') as zf:
        for filename in zf.namelist():
            # Avoid processing system files or directories
            if not filename.startswith('__MACOSX/') and not filename.endswith('/'):
                full_text.append(f"--- Zipped File: {filename} ---\n")
                with zf.open(filename) as unzipped_file:
                    unzipped_stream = io.BytesIO(unzipped_file.read())
                    # Recursively call the main converter function
                    converted_content = universal_file_converter(filename, unzipped_stream)
                    full_text.append(converted_content)
                    full_text.append("\n" + "="*40 + "\n")
    return "\n".join(full_text)

# --- Main Controller Function ---

def universal_file_converter(filename, file_stream):
    """Acts as a router to call the correct conversion function based on file extension."""
    _, extension = os.path.splitext(filename.lower())

    # Route to the appropriate function
    if extension == '.docx':
        return convert_docx_to_text(file_stream)
    elif extension == '.xlsx':
        return convert_xlsx_to_text(file_stream)
    elif extension == '.pptx':
        return convert_pptx_to_text(file_stream)
    elif extension in ['.html', '.htm']:
        return convert_html_to_markdown(file_stream)
    elif extension == '.zip':
        return process_zip_file(file_stream)
    else:
        # Fallback for plain text or unsupported files
        try:
            return file_stream.read().decode('utf-8')
        except Exception:
            return f"Unsupported file type: {extension}. Could not read file."

# --- Streamlit App UI ---

# Set the page title and a brief description
st.set_page_config(page_title="File-to-Text Converter", layout="centered")
st.title("📄 Universal File-to-Text Converter")
st.markdown("Drag, drop, and download. It's that simple.")
st.markdown("Supports: `.docx`, `.xlsx`, `.pptx`, `.html`, `.zip`")

# [1] File uploader widget for user input
uploaded_file = st.file_uploader(
    "Drag and drop your file here",
    type=['docx', 'xlsx', 'pptx', 'html', 'htm', 'zip'],
    label_visibility="collapsed"
)

if uploaded_file is not None:
    # Get filename and prepare a file stream
    filename = uploaded_file.name
    file_stream = io.BytesIO(uploaded_file.getvalue())

    st.info(f"Processing `{filename}`...")

    # [2] Process the file into a text string
    full_text = universal_file_converter(filename, file_stream)

    # [3] Display a preview of the first 1000 characters
    st.subheader("✅ Conversion Successful! Here's a Preview:")
    st.text_area(
        "Showing first 1000 characters of the output",
        full_text[:1000],
        height=250
    )

    # [4] Offer the full text file for download
    output_filename = os.path.splitext(filename)[0] + '.txt'
    st.download_button(
        label="📥 Download Full Text File",
        data=full_text.encode('utf-8'),
        file_name=output_filename,
        mime='text/plain',
    )
